import streamlit as st
import time
import os
from agents import (
    market_research,
    competitor_analysis,
    customer_insight,
    swot_analysis,
    financial_estimation,
    risk_feasibility,
    final_strategy
)
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import textwrap
from pptx import Presentation

# Set page config
st.set_page_config(
    page_title="Business Strategy AI Agent",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for styling
st.markdown("""
    <style>
    .main {
        background-color: #0e1117;
    }
    .stButton>button {
        width: 100%;
        background-color: #ff4b4b;
        color: white;
        border-radius: 5px;
        height: 3em;
        font-weight: bold;
    }
    .stTextInput>div>div>input {
        background-color: #262730;
        color: white;
    }
    h1 {
        color: #ff4b4b;
    }
    h2, h3 {
        color: #fafafa;
    }
    .report-text {
        color: #c9c9c9;
        font-family: 'Courier New', Courier, monospace;
    }
    </style>
    """, unsafe_allow_html=True)

def generate_pdf(filename, content):
    try:
        pdf_file = filename.replace(".md", ".pdf")
        c = canvas.Canvas(pdf_file, pagesize=letter)
        width, height = letter
        text = c.beginText(40, height - 40)
        text.setFont("Helvetica", 10)
        
        for line in content.split('\n'):
            wrapped_lines = textwrap.wrap(line, width=90)
            for w_line in wrapped_lines:
                if text.getY() < 40:
                    c.drawText(text)
                    c.showPage()
                    text = c.beginText(40, height - 40)
                    text.setFont("Helvetica", 10)
                text.textLine(w_line)
        c.drawText(text)
        c.save()
        return pdf_file
    except Exception as e:
        return None

def generate_ppt(filename, business_idea, content_dict):
    try:
        ppt_file = filename.replace(".md", ".pptx")
        prs = Presentation()
        
        # Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Business Strategy Report"
        subtitle.text = business_idea

        def add_slide(title_text, content_text):
            layout = prs.slide_layouts[1] 
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            tf = slide.placeholders[1].text_frame
            tf.text = content_text[:1000] + "..." if len(content_text) > 1000 else content_text

        for section, text in content_dict.items():
            add_slide(section, text)

        prs.save(ppt_file)
        return ppt_file
    except Exception as e:
        return None

# Sidebar
with st.sidebar:
    st.image("https://cdn-icons-png.flaticon.com/512/2083/2083256.png", width=50)
    st.title("AI Strategy Agent")
    st.markdown("---")
    business_idea = st.text_input("💡 Business Idea", placeholder="e.g. AI-powered software agency")
    industry = st.text_input("🏭 Industry", placeholder="e.g. SaaS")
    region = st.text_input("🌍 Target Region", placeholder="e.g. India, Jaipur")
    
    st.markdown("### ⚙️ Settings")
    api_status = "✅ Connected" if os.getenv("GROQ_API_KEY") or os.getenv("GEMINI_API_KEY") else "❌ No Keys Found"
    st.info(f"API Status: {api_status}")
    
    run_btn = st.button("🚀 Generate Strategy")

# Main Content
st.title("Business Strategy & Market Research System")
st.markdown("#### From Idea to Execution Plan in Minutes")

if run_btn:
    if not business_idea or not industry or not region:
        st.error("Please fill in all fields (Idea, Industry, Region).")
    else:
        results = {}
        full_report_text = ""
        
        # Container for progress
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        # Tabs for live viewing
        tab1, tab2, tab3, tab4, tab5, tab6, tab7 = st.tabs([
            "📊 Market", "⚔️ Competitors", "👥 Customers", "🛡️ SWOT", "💰 Financials", "⚠️ Risks", "🏁 Final Strategy"
        ])

        # 1. Market Research
        status_text.markdown("### 🔍 Agent 1/7: Analyzing Market...")
        market_res = market_research.run(business_idea, industry, region)
        results["Market Research"] = market_res
        with tab1:
            st.markdown(market_res)
        progress_bar.progress(15)
        time.sleep(2)

        # 2. Competitor Analysis
        status_text.markdown("### 🕵️ Agent 2/7: Scouting Competitors...")
        comp_analysis = competitor_analysis.run(business_idea, industry)
        results["Competitor Analysis"] = comp_analysis
        with tab2:
            st.markdown(comp_analysis)
        progress_bar.progress(30)
        time.sleep(2)

        # 3. Customer Insight
        status_text.markdown("### 🧠 Agent 3/7: Understanding Customers...")
        cust_insights = customer_insight.run(business_idea, industry)
        results["Customer Insights"] = cust_insights
        with tab3:
            st.markdown(cust_insights)
        progress_bar.progress(45)
        time.sleep(2)

        # 4. SWOT
        status_text.markdown("### ⚖️ Agent 4/7: Developing SWOT...")
        swot_out = swot_analysis.run(business_idea, market_res, comp_analysis, cust_insights)
        results["SWOT Analysis"] = swot_out
        with tab4:
            st.markdown(swot_out)
        progress_bar.progress(60)
        time.sleep(2)

        # 5. Financial
        status_text.markdown("### 💸 Agent 5/7: Estimating Financials...")
        fin_est = financial_estimation.run(business_idea, industry, region)
        results["Financial Estimation"] = fin_est
        with tab5:
            st.markdown(fin_est)
        progress_bar.progress(75)
        time.sleep(2)

        # 6. Risk
        status_text.markdown("### 🚧 Agent 6/7: Assessing Risks...")
        risk_out = risk_feasibility.run(business_idea, industry, region)
        results["Risk Assessment"] = risk_out
        with tab6:
            st.markdown(risk_out)
        progress_bar.progress(90)
        time.sleep(2)

        # 7. Final Strategy
        status_text.markdown("### 🎯 Agent 7/7: Formulating Final Strategy...")
        final_strat = final_strategy.run(
            business_idea, market_res, comp_analysis, cust_insights, swot_out, fin_est, risk_out
        )
        results["Final Strategy"] = final_strat
        with tab7:
            st.markdown(final_strat)
        progress_bar.progress(100)
        status_text.success("✅ Analysis Complete!")

        # Compile Full Report
        full_report_text = f"# Business Strategy Report: {business_idea}\n\n"
        for key, val in results.items():
            full_report_text += f"## {key}\n{val}\n\n"

        # Generate Files
        filename = f"Strategy_Report_{business_idea.replace(' ', '_')}.md"
        with open(filename, "w") as f:
            f.write(full_report_text)
            
        pdf_path = generate_pdf(filename, full_report_text)
        ppt_path = generate_ppt(filename, business_idea, results)
        
        st.markdown("---")
        st.subheader("📥 Download Report")
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.download_button("📄 Download Markdown", full_report_text, file_name=filename)
        
        with col2:
            if pdf_path:
                with open(pdf_path, "rb") as pdf_file:
                    st.download_button("📕 Download PDF", pdf_file, file_name=os.path.basename(pdf_path))
        
        with col3:
            if ppt_path:
                with open(ppt_path, "rb") as ppt_file:
                    st.download_button("📊 Download PPT", ppt_file, file_name=os.path.basename(ppt_path))

else:
    st.info("👈 Enter your business details in the sidebar and click Generate!")
