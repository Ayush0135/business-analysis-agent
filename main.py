import os
import sys

# Add current directory to path so imports work
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from agents import (
    market_research,
    competitor_analysis,
    customer_insight,
    swot_analysis,
    financial_estimation,
    risk_feasibility,
    final_strategy
)

def main():
    print("Welcome to the Business Strategy & Market Research Multi-Agent System")
    print("---------------------------------------------------------------------")
    
    business_idea = input("Enter your Business Idea: ").strip()
    industry = input("Enter the Industry (e.g., EdTech, SaaS, F&B): ").strip()
    region = input("Enter Target Region (e.g., India, USA, Global): ").strip()
    
    if not business_idea:
        print("Business idea is required.")
        return

    import time

    # 1. Market Research
    market_res = market_research.run(business_idea, industry, region)
    time.sleep(5) # Polite delay
    
    # 2. Competitor Analysis
    comp_analysis = competitor_analysis.run(business_idea, industry)
    time.sleep(5) # Polite delay
    
    # 3. Customer Insight
    cust_insights = customer_insight.run(business_idea, industry)
    time.sleep(5) # Polite delay
    
    # 4. SWOT
    swot_out = swot_analysis.run(business_idea, market_res, comp_analysis, cust_insights)
    time.sleep(5) # Polite delay
    
    # 5. Financial
    fin_est = financial_estimation.run(business_idea, industry, region)
    time.sleep(5) # Polite delay
    
    # 6. Risk
    risk_out = risk_feasibility.run(business_idea, industry, region)
    time.sleep(5) # Polite delay
    
    # 7. Final Strategy
    final_strat = final_strategy.run(
        business_idea, 
        market_res, 
        comp_analysis, 
        cust_insights, 
        swot_out, 
        fin_est, 
        risk_out
    )
    
    # Combine everything into a report
    full_report = f"""
# Business Strategy Report: {business_idea}

## 1. Market Research
{market_res}

## 2. Competitor Analysis
{comp_analysis}

## 3. Customer Insights
{cust_insights}

## 4. SWOT Analysis
{swot_out}

## 5. Financial Estimation
{fin_est}

## 6. Risk & Feasibility
{risk_out}

## 7. Final Strategy & Verdict
{final_strat}
    """
    
    filename = f"Strategy_Report_{business_idea.replace(' ', '_')}.md"
    with open(filename, "w") as f:
        f.write(full_report)
        
    print(f"\n\nReport generated successfully: {filename}")
    
    # Attempt PDF generation (simplified)
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        import textwrap
        
        pdf_file = filename.replace(".md", ".pdf")
        c = canvas.Canvas(pdf_file, pagesize=letter)
        width, height = letter
        
        text = c.beginText(40, height - 40)
        text.setFont("Helvetica", 10)
        
        # Very basic split by lines
        for line in full_report.split('\n'):
            wrapped_lines = textwrap.wrap(line, width=90)
            for w_line in wrapped_lines:
                if text.getY() < 40:
                    c.drawText(text)
                    c.showPage()
                    text = c.beginText(40, height - 40)
                    text.setFont("Helvetica", 10)
                text.textLine(w_line)
                
        c.drawText(text)
        c.save()
        print(f"PDF generated: {pdf_file}")
        
    except Exception as e:
        print(f"Could not generate PDF: {e}")

    # Attempt PPT generation
    try:
        from pptx import Presentation
        from pptx.util import Inches

        ppt_file = filename.replace(".md", ".pptx")
        prs = Presentation()
        
        # Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Business Strategy Report"
        subtitle.text = business_idea

        def add_slide(title_text, content_text):
            layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            # Basic content placement - handling large text is tricky in PPT
            # Truncate or just place it
            tf = slide.placeholders[1].text_frame
            tf.text = content_text[:1000] + "..." if len(content_text) > 1000 else content_text

        add_slide("Market Research", market_res)
        add_slide("Competitor Analysis", comp_analysis)
        add_slide("Customer Insights", cust_insights)
        add_slide("SWOT Analysis", swot_out)
        add_slide("Financial Estimation", fin_est)
        add_slide("Risk Assessment", risk_out)
        add_slide("Final Strategy", final_strat)

        prs.save(ppt_file)
        print(f"PPT generated: {ppt_file}")

    except Exception as e:
        print(f"Could not generate PPT: {e}")

if __name__ == "__main__":
    main()
