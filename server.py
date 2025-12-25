from flask import Flask, render_template, request, Response, jsonify
from flask_sqlalchemy import SQLAlchemy
from datetime import datetime
import json
import time
import os
import shutil
from agents import (
    market_research,
    competitor_analysis,
    customer_insight,
    swot_analysis,
    financial_estimation,
    risk_feasibility,
    final_strategy
)
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import textwrap
from pptx import Presentation

app = Flask(__name__)

# Database Configuration
# Use DATABASE_URL from environment (Vercel/Cloud) or fallback to local SQLite
database_url = os.getenv('DATABASE_URL', 'sqlite:///reports.db')
if database_url.startswith("postgres://"):
    database_url = database_url.replace("postgres://", "postgresql://", 1)

app.config['SQLALCHEMY_DATABASE_URI'] = database_url
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)

# --- Database Model ---
class Analysis(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    idea = db.Column(db.String(200), nullable=False)
    industry = db.Column(db.String(100), nullable=False)
    region = db.Column(db.String(100), nullable=False)
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)
    md_path = db.Column(db.String(200))
    pdf_path = db.Column(db.String(200))
    ppt_path = db.Column(db.String(200))

    def to_dict(self):
        return {
            "id": self.id,
            "idea": self.idea,
            "industry": self.industry,
            "region": self.region,
            "date": self.timestamp.strftime("%Y-%m-%d %H:%M"),
            "files": {
                "md": self.md_path,
                "pdf": self.pdf_path,
                "ppt": self.ppt_path
            }
        }

# Create DB Tables
with app.app_context():
    db.create_all()

# --- Helper Functions for File Generation ---
def create_pdf(filename, content):
    try:
        c = canvas.Canvas(filename, pagesize=letter)
        width, height = letter
        text = c.beginText(40, height - 40)
        text.setFont("Helvetica", 10)
        for line in content.split('\n'):
            wrapped = textwrap.wrap(line, width=90)
            for w in wrapped:
                if text.getY() < 40:
                    c.drawText(text)
                    c.showPage()
                    text = c.beginText(40, height - 40)
                    text.setFont("Helvetica", 10)
                text.textLine(w)
        c.drawText(text)
        c.save()
        return filename
    except: return None

def create_ppt(filename, idea, data):
    try:
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Business Strategy Report"
        title_slide.placeholders[1].text = idea
        
        for section, content in data.items():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section
            tf = slide.placeholders[1].text_frame
            tf.text = content[:1000] + "..." if len(content) > 1000 else content
            
        prs.save(filename)
        return filename
    except: return None

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/api/history')
def get_history():
    reports = Analysis.query.order_by(Analysis.timestamp.desc()).all()
    return jsonify([r.to_dict() for r in reports])

@app.route('/stream_analysis')
def stream_analysis():
    idea = request.args.get('idea')
    industry = request.args.get('industry')
    region = request.args.get('region')

    if not idea or not industry or not region:
        return Response("data: Error: Missing fields\n\n", mimetype='text/event-stream')

    def generate():
        with app.app_context(): # Need app context for DB operations if we were writing immediately, but we write at end
            results = {}
            
            # 1. Market Research
            yield f"data: {json.dumps({'step': 1, 'name': 'Market Research', 'status': 'running'})}\n\n"
            market_res = market_research.run(idea, industry, region)
            results['Market Research'] = market_res
            yield f"data: {json.dumps({'step': 1, 'name': 'Market Research', 'status': 'done', 'content': market_res})}\n\n"
            
            # 2. Competitor
            yield f"data: {json.dumps({'step': 2, 'name': 'Competitor Analysis', 'status': 'running'})}\n\n"
            comp_res = competitor_analysis.run(idea, industry)
            results['Competitor Analysis'] = comp_res
            yield f"data: {json.dumps({'step': 2, 'name': 'Competitor Analysis', 'status': 'done', 'content': comp_res})}\n\n"

            # 3. Customer
            yield f"data: {json.dumps({'step': 3, 'name': 'Customer Insights', 'status': 'running'})}\n\n"
            cust_res = customer_insight.run(idea, industry)
            results['Customer Insights'] = cust_res
            yield f"data: {json.dumps({'step': 3, 'name': 'Customer Insights', 'status': 'done', 'content': cust_res})}\n\n"

            # 4. SWOT
            yield f"data: {json.dumps({'step': 4, 'name': 'SWOT Analysis', 'status': 'running'})}\n\n"
            swot_res = swot_analysis.run(idea, market_res, comp_res, cust_res)
            results['SWOT Analysis'] = swot_res
            yield f"data: {json.dumps({'step': 4, 'name': 'SWOT Analysis', 'status': 'done', 'content': swot_res})}\n\n"

            # 5. Financial
            yield f"data: {json.dumps({'step': 5, 'name': 'Financial Estimation', 'status': 'running'})}\n\n"
            fin_res = financial_estimation.run(idea, industry, region)
            results['Financial Estimation'] = fin_res
            yield f"data: {json.dumps({'step': 5, 'name': 'Financial Estimation', 'status': 'done', 'content': fin_res})}\n\n"

            # 6. Risk
            yield f"data: {json.dumps({'step': 6, 'name': 'Risk Assessment', 'status': 'running'})}\n\n"
            risk_res = risk_feasibility.run(idea, industry, region)
            results['Risk Assessment'] = risk_res
            yield f"data: {json.dumps({'step': 6, 'name': 'Risk Assessment', 'status': 'done', 'content': risk_res})}\n\n"

            # 7. Final Strategy
            yield f"data: {json.dumps({'step': 7, 'name': 'Final Strategy', 'status': 'running'})}\n\n"
            final_res = final_strategy.run(idea, market_res, comp_res, cust_res, swot_res, fin_res, risk_res)
            results['Final Strategy'] = final_res
            yield f"data: {json.dumps({'step': 7, 'name': 'Final Strategy', 'status': 'done', 'content': final_res})}\n\n"

            # Generate Files & Save to DB
            yield f"data: {json.dumps({'step': 8, 'name': 'Generating Files', 'status': 'running'})}\n\n"
            
            # Create a dedicated directory for reports if needed, but static/reports is fine
            report_dir = "static/reports"
            if not os.path.exists(report_dir):
                os.makedirs(report_dir)

            timestamp = int(time.time())
            base_name = f"{idea.replace(' ', '_')}_{timestamp}"
            
            md_path = f"{report_dir}/{base_name}.md"
            pdf_path = f"{report_dir}/{base_name}.pdf"
            ppt_path = f"{report_dir}/{base_name}.pptx"
            
            full_text = f"# Business Strategy: {idea}\n\n" + "\n\n".join([f"## {k}\n{v}" for k,v in results.items()])
            
            with open(md_path, "w") as f: f.write(full_text)
            create_pdf(pdf_path, full_text)
            create_ppt(ppt_path, idea, results)
            
            # Save to Database
            new_analysis = Analysis(
                idea=idea,
                industry=industry,
                region=region,
                md_path=md_path,
                pdf_path=pdf_path,
                ppt_path=ppt_path
            )
            db.session.add(new_analysis)
            db.session.commit()
            
            yield f"data: {json.dumps({'step': 8, 'name': 'Complete', 'status': 'complete', 'files': {'md': md_path, 'pdf': pdf_path, 'ppt': ppt_path}})}\n\n"

    return Response(generate(), mimetype='text/event-stream')

if __name__ == '__main__':
    # Ensure static exists for reports
    if not os.path.exists('static'):
        os.makedirs('static')
    app.run(debug=True, port=3000)
